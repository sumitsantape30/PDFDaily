# from rest_framework.decorators import api_view
# from rest_framework.response import Response
# from pptx import Presentation
# from docx import Document
# import tempfile
# import comtypes.client
# import os
# import pythoncom


# @api_view(['POST'])
# def ppt_to_word(request):
#     pythoncom.CoInitialize()
#     ppt_file = request.FILES.get('ppt_file')
    
#     if not ppt_file.name.endswith(('.ppt', '.pptx')):
#         return Response({'error': 'Invalid file type'})

   
#     with tempfile.NamedTemporaryFile(delete=False) as ppt_path:
#         # ppt_path = save_temporary_file(ppt_file)
#         docx_path = convert_ppt_to_docx(ppt_path)
#         docx_file = create_file_response(docx_path)
#         return docx_file
#     try:
#         return Response({'message': 'No PPT file provided.'}, status=400)

        
#         docx_file = 'converted.docx'
#         doc = Document()
        
#         for slide in presentation.slides:
#             for shape in slide.shapes:
#                 if shape.has_text_frame:
#                     for paragraph in shape.text_frame.paragraphs:
#                         doc.add_paragraph(paragraph.text)
        
#         doc.save(docx_file)
        
#         return Response({'message': 'Conversion successful!', 'docx_file': docx_file})
#     finally:
#         # Uninitialize the COM library
#         pythoncom.CoUninitialize()

# import os
# import tempfile
# import comtypes.client
# from docx import Document
# from rest_framework.decorators import api_view
# from rest_framework.response import Response

# @api_view(['POST'])
# def ppt_to_word(request):
#     ppt_file = request.FILES.get('ppt_file')

#     if ppt_file:
#         comtypes.CoInitialize()  # Initialize COM environment

#         ppt_path = save_temporary_file(ppt_file)
#         docx_path = convert_ppt_to_docx(ppt_path)
#         docx_file = create_file_response(docx_path)
#         return docx_file
#     else:
#         return Response({'message': 'No PPT file provided.'}, status=400)

# def save_temporary_file(file):
#     temp_dir = tempfile.gettempdir()
#     file_path = os.path.join(temp_dir, file.name)
#     with open(file_path, 'wb') as temp_file:
#         for chunk in file.chunks():
#             temp_file.write(chunk)
#     return file_path

# # def convert_ppt_to_docx(ppt_path):
# #     powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
# #     powerpoint.Visible = 1

# #     presentation = powerpoint.Presentations.Open(ppt_path)
# #     docx_path = ppt_path + '.docx'

# #     presentation.SaveAs(docx_path, 16)  # 16 represents the file format for DOCX
# #     presentation.Close()
# #     powerpoint.Quit()

# #     return docx_path
# def convert_ppt_to_docx(ppt_path):
#     powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
#     powerpoint.Visible = 1

#     presentation = powerpoint.Presentations.Open(ppt_path)
#     docx_path = ppt_path + '.docx'
#     print(docx_path)  # Add this line to check the generated path

#     presentation.SaveAs(docx_path,16)  # 16 represents the file format for DOCX
#     presentation.Close()
#     powerpoint.Quit()

#     return docx_path

# def create_file_response(file_path):
#     with open(file_path, 'rb') as docx_file:
#         response = Response(docx_file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
#         response['Content-Disposition'] = 'attachment; filename=' + os.path.basename(file_path)
#         return response

from pptx import Presentation
from docx import Document
from rest_framework.decorators import api_view
from rest_framework.response import Response
import os

@api_view(['POST'])
def ppt_to_word(request):
    ppt_file = request.FILES.get('ppt_file')

    if ppt_file:
        ppt_path = save_temporary_file(ppt_file)
        docx_file = convert_ppt_to_docx(ppt_path)
        return docx_file
    else:
        return Response({'message': 'No PPT file provided.'}, status=400)


    file_path = os.path.join(temp_dir, file.name)
    with open(file_path, 'wb') as temp_file:
        for chunk in file.chunks():
            temp_file.write(chunk)
    return file_path


    presentation = Presentation(ppt_path)
    doc = Document()

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    doc.add_paragraph(paragraph.text)

    docx_file = 'converted.docx'
    doc.save(docx_file)

    with open(docx_file, 'rb') as file:
        response = Response(file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        response['Content-Disposition'] = 'attachment; filename=converted.docx'
        return response
